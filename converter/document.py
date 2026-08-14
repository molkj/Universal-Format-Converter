"""文档格式转换：PDF / Word / Excel / PPT / TXT / MD / HTML"""
from __future__ import annotations

import csv
import html.parser
import io
import os

from .utils import (
    ConverterError,
    Progress,
    ensure_dir,
    get_ext,
    stem,
    unique_path,
)

# ---------------------------------------------------------------------------
# 内部工具
# ---------------------------------------------------------------------------

class _TextExtractor(html.parser.HTMLParser):
    """从 HTML 中提取纯文本"""

    def __init__(self):
        super().__init__()
        self._parts: list[str] = []
        self._skip = 0

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style"):
            self._skip += 1
        if tag in ("p", "br", "div", "li", "h1", "h2", "h3", "h4", "tr"):
            self._parts.append("\n")

    def handle_endtag(self, tag):
        if tag in ("script", "style") and self._skip > 0:
            self._skip -= 1

    def handle_data(self, data):
        if not self._skip:
            self._parts.append(data)

    def text(self) -> str:
        lines = ["".join(self._parts).strip()]
        return "\n".join(l for l in lines)


def _html_to_text(src: str) -> str:
    with open(src, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    p = _TextExtractor()
    p.feed(content)
    return p.text()


def _extract_docx_text(src: str) -> str:
    from docx import Document

    doc = Document(src)
    parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    return "\n".join(parts)


def _write_docx_from_text(text: str, out_path: str):
    from docx import Document

    doc = Document()
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("# "):
            doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith("### "):
            doc.add_heading(stripped[4:], level=3)
        elif stripped.startswith(("- ", "* ", "+ ")):
            doc.add_paragraph(stripped[2:], style="List Bullet")
        else:
            doc.add_paragraph(stripped)
    doc.save(out_path)


def _office_com_convert(src: str, out_path: str, app_ids, save_format, load_format=None):
    """通过 Windows COM 调用 Office/WPS 转换（如 Word→PDF、PPT→PDF）"""
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        raise ConverterError(
            "未找到本机 Office/WPS 的 COM 接口（需要安装 Microsoft Office 或 WPS，"
            "并且 Python 需包含 pywin32）。"
        )

    pythoncom.CoInitialize()
    app = None
    last_err = None
    for prog_id in app_ids:
        try:
            app = win32com.client.Dispatch(prog_id)
            break
        except Exception as e:  # noqa: BLE001
            last_err = e
    if app is None:
        pythoncom.CoUninitialize()
        raise ConverterError(
            f"未检测到可用的 Office/WPS 程序（尝试了 {', '.join(app_ids)}）。\n"
            "请安装 Microsoft Office 或 WPS Office 后重试。"
        )
    try:
        app.Visible = False
        app.DisplayAlerts = False
        doc = app.Documents.Open(os.path.abspath(src), ReadOnly=True)
        try:
            doc.SaveAs(os.path.abspath(out_path), FileFormat=save_format)
        finally:
            doc.Close(False)
    except Exception as e:  # noqa: BLE001
        last_err = e
        raise ConverterError(f"Office 转换失败：{e}")
    finally:
        try:
            app.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()
    return out_path


# ---------------------------------------------------------------------------
# PDF 相关
# ---------------------------------------------------------------------------

def pdf_to_docx(src, out_path, progress):
    from pdf2docx import Converter

    progress.report(10, 100, "正在解析 PDF…")
    cv = Converter(src)
    try:
        cv.convert(out_path)
    finally:
        cv.close()
    progress.report(100, 100, "PDF → Word 完成")
    return out_path


def pdf_to_text(src, out_path, progress):
    import fitz

    doc = fitz.open(src)
    parts = []
    total = doc.page_count
    for i, page in enumerate(doc):
        if progress.cancelled:
            raise InterruptedError
        parts.append(page.get_text("text"))
        progress.report(i + 1, total, f"提取第 {i + 1}/{total} 页文本…")
    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n\n".join(parts))
    progress.report(100, 100, "PDF → 文本 完成")
    return out_path


def pdf_to_images(src, out_path, dpi=150, image_format="png", progress=None):
    """PDF 每页导出为一张图片，返回 [文件1, 文件2, ...]"""
    import fitz

    fmt = image_format.lower()
    ext_map = {"jpg": "jpg", "jpeg": "jpg", "png": "png", "webp": "webp"}
    ext = ext_map.get(fmt, "png")
    matrix = fitz.Matrix(dpi / 72, dpi / 72)
    doc = fitz.open(src)
    total = doc.page_count
    results = []
    for i, page in enumerate(doc):
        if progress and progress.cancelled:
            raise InterruptedError
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        out_file = unique_path(out_path, f"{stem(src)}_第{i + 1}页", ext)
        if fmt == "jpg":
            pix.pil_save(out_file, format="JPEG", quality=90)
        else:
            pix.pil_save(out_file, format="PNG" if ext == "png" else "WEBP")
        results.append(out_file)
        if progress:
            progress.report(i + 1, total, f"导出第 {i + 1}/{total} 页…")
    return results


# ---------------------------------------------------------------------------
# Office 相关
# ---------------------------------------------------------------------------

def docx_to_pdf(src, out_path, progress):
    progress.report(10, 100, "正在调用 Office/WPS…")
    return _office_com_convert(
        src, out_path,
        app_ids=("Word.Application", "KWps.Application", "WPS.Application"),
        save_format=17,  # wdFormatPDF
    )


def docx_to_text(src, out_path, progress):
    text = _extract_docx_text(src)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(text)
    progress.report(100, 100, "Word → 文本 完成")
    return out_path


def docx_to_md(src, out_path, progress):
    text = _extract_docx_text(src)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(text)
    progress.report(100, 100, "Word → Markdown 完成")
    return out_path


def text_to_docx(src, out_path, progress):
    with open(src, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    _write_docx_from_text(text, out_path)
    progress.report(100, 100, "文本 → Word 完成")
    return out_path


def markdown_to_docx(src, out_path, progress):
    with open(src, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    _write_docx_from_text(text, out_path)
    progress.report(100, 100, "Markdown → Word 完成")
    return out_path


def txt_to_md(src, out_path, progress):
    """纯文本 → Markdown：逐行转为简单 Markdown 结构"""
    with open(src, "r", encoding="utf-8", errors="replace") as f:
        lines = f.read().splitlines()
    parts = []
    for line in lines:
        s = line.strip()
        if not s:
            continue
        if s.startswith("### "):
            parts.append(f"### {s[4:]}")
        elif s.startswith("## "):
            parts.append(f"## {s[3:]}")
        elif s.startswith("# "):
            parts.append(f"# {s[2:]}")
        elif s.startswith(("- ", "* ", "+ ")):
            parts.append(f"- {s[2:]}")
        else:
            parts.append(s)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n\n".join(parts))
    progress.report(100, 100, "文本 → Markdown 完成")
    return out_path


def html_to_text(src, out_path, progress):
    text = _html_to_text(src)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(text)
    progress.report(100, 100, "HTML → 文本 完成")
    return out_path


def html_to_docx(src, out_path, progress):
    text = _html_to_text(src)
    _write_docx_from_text(text, out_path)
    progress.report(100, 100, "HTML → Word 完成")
    return out_path


def md_to_html(src, out_path, progress):
    try:
        import markdown
    except ImportError:
        # 极简降级：逐行渲染标题/粗体/斜体/代码
        with open(src, "r", encoding="utf-8", errors="replace") as f:
            lines = f.read().splitlines()
        body = []
        for line in lines:
            s = line.strip()
            if s.startswith("### "):
                body.append(f"<h3>{s[4:]}</h3>")
            elif s.startswith("## "):
                body.append(f"<h2>{s[3:]}</h2>")
            elif s.startswith("# "):
                body.append(f"<h1>{s[2:]}</h1>")
            elif s.startswith(("- ", "* ")):
                body.append(f"<li>{s[2:]}</li>")
            elif s == "---":
                body.append("<hr>")
            elif s:
                body.append(f"<p>{s}</p>")
        body_html = "\n".join(body)
    else:
        body_html = markdown.markdown(
            open(src, "r", encoding="utf-8", errors="replace").read(),
            extensions=["extra"],
        )
    html_doc = (
        "<!DOCTYPE html><html lang='zh'><head><meta charset='utf-8'>"
        f"<title>{stem(src)}</title></head><body>{body_html}</body></html>"
    )
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(html_doc)
    progress.report(100, 100, "Markdown → HTML 完成")
    return out_path


# ---------------------------------------------------------------------------
# Excel 相关
# ---------------------------------------------------------------------------

def xlsx_to_csv(src, out_path, progress):
    """导出第一个工作表为 CSV"""
    from openpyxl import load_workbook

    wb = load_workbook(src, read_only=True, data_only=True)
    ws = wb.worksheets[0]
    with open(out_path, "w", encoding="utf-8-sig", newline="") as f:
        writer = csv.writer(f)
        for row in ws.iter_rows(values_only=True):
            writer.writerow(["" if c is None else c for c in row])
    wb.close()
    progress.report(100, 100, "Excel → CSV 完成")
    return out_path


def xlsx_to_txt(src, out_path, progress):
    from openpyxl import load_workbook

    wb = load_workbook(src, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"【工作表：{ws.title}】")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c).strip() for c in row]
            if any(cells):
                parts.append("\t".join(cells))
    wb.close()
    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n".join(parts))
    progress.report(100, 100, "Excel → 文本 完成")
    return out_path


def xls_to_xlsx(src, out_path, progress):
    try:
        import xlrd
    except ImportError:
        raise ConverterError("转换 .xls 需要 xlrd 库支持。")
    try:
        from openpyxl import Workbook
    except ImportError:
        raise ConverterError("缺少 openpyxl 库。")

    wb_old = xlrd.open_workbook(src)
    wb = Workbook()
    # 兼容 xlrd 1.x（.sheets()）与 xlrd 2.x（.sheet_names() / .sheet_by_name()）
    if hasattr(wb_old, "sheets"):
        sheets = wb_old.sheets()
    else:
        sheets = [wb_old.sheet_by_name(n) for n in wb_old.sheet_names()]
    for sheet in sheets:
        ws = wb.create_sheet(title=sheet.name[:31])
        for r in range(sheet.nrows):
            ws.append(sheet.row_values(r))
    if "Sheet" in wb.sheetnames:
        del wb["Sheet"]
    wb.save(out_path)
    progress.report(100, 100, "XLS → XLSX 完成")
    return out_path


def csv_to_xlsx(src, out_path, progress):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    with open(src, "r", encoding="utf-8-sig", errors="replace", newline="") as f:
        reader = csv.reader(f)
        for row in reader:
            ws.append(row)
    wb.save(out_path)
    progress.report(100, 100, "CSV → Excel 完成")
    return out_path


def pptx_to_pdf(src, out_path, progress):
    progress.report(10, 100, "正在调用 Office/WPS…")
    return _office_com_convert(
        src, out_path,
        app_ids=("PowerPoint.Application", "KWpp.Application", "WPP.Application"),
        save_format=32,  # ppSaveAsPDF
    )


def pptx_to_text(src, out_path, progress):
    from pptx import Presentation

    prs = Presentation(src)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"=== 第 {i} 页 ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        parts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append("\t".join(c.text.strip() for c in row.cells))
    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n".join(parts))
    progress.report(100, 100, "PPT → 文本 完成")
    return out_path


# ---------------------------------------------------------------------------
# 注册
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# 文本编码转换（GBK ↔ UTF-8，修复乱码）
# ---------------------------------------------------------------------------

def detect_encoding(src: str) -> str:
    """检测文本文件编码：优先 BOM，其次尝试 UTF-8/GBK"""
    with open(src, "rb") as f:
        raw = f.read(4096)
    if raw.startswith(b"\xef\xbb\xbf"):
        return "utf-8-sig"
    if raw.startswith(b"\xff\xfe") or raw.startswith(b"\xfe\xff"):
        return "utf-16"
    for enc in ("utf-8", "gbk"):
        try:
            raw.decode(enc)
            return enc
        except UnicodeDecodeError:
            continue
    return "gbk"  # 兜底（GBK 几乎能解码任意字节）


def convert_text_encoding(src: str, out_path: str, target_enc: str,
                          progress) -> str:
    """把文本文件转换为目标编码（utf-8 / utf-8-sig / gbk）

    目标编码无法表示的字符（如 GBK 遇 emoji）自动替换为「?」，
    避免转换崩溃；进度消息里会提示替换数量。"""
    src_enc = detect_encoding(src)
    with open(src, "rb") as f:
        raw = f.read()
    try:
        text = raw.decode(src_enc, errors="replace")
    except Exception as e:  # noqa: BLE001
        raise ConverterError(f"读取文件编码失败（{src_enc}）：{e}") from e
    # 统计无法用目标编码表示的字符数（如 GBK 不支持的 emoji）
    replaced = 0
    for ch in text:
        try:
            ch.encode(target_enc)
        except UnicodeEncodeError:
            replaced += 1
    # 写入目标编码；不支持的字符替换为 ?（errors=replace 不崩溃）
    with open(out_path, "w", encoding=target_enc, newline="",
              errors="replace") as f:
        f.write(text)
    progress.report(100, 100,
                    f"编码转换完成（{src_enc} → {target_enc}"
                    + (f"，{replaced} 个字符无法表示已替换为 ?" if replaced else "")
                    + "）")
    return out_path


def _encoding_wrapper(target_enc: str, label: str):
    def wrapper(src, tgt, out_dir, progress, opts):
        # 输出名带编码标记，避免与普通 txt 冲突（如 报告_UTF8.txt）
        src_enc = detect_encoding(src)
        if src_enc == target_enc or (src_enc == "utf-8-sig" and target_enc == "utf-8"):
            raise ConverterError("文件已是目标编码，无需转换。")
        marker = "UTF8" if target_enc in ("utf-8", "utf-8-sig") else "GBK"
        out = unique_path(out_dir, f"{stem(src)}_{marker}", "txt")
        progress.report(0, 100, f"开始：{label}")
        result = convert_text_encoding(src, out, target_enc, progress)
        return [result]

    return wrapper


def register(registry):
    """向全局注册表注册本模块支持的 (源格式, 目标格式) -> 转换函数"""

    # 每种转换的目标扩展名列表（用于 GUI 展示）
    targets = {
        "pdf": ["docx", "txt", "png", "jpg", "webp"],
        "docx": ["pdf", "txt", "md"],
        "doc": ["pdf", "txt", "docx"],
        "xlsx": ["csv", "txt"],
        "xls": ["xlsx"],
        "csv": ["xlsx", "utf8", "gbk"],
        "txt": ["docx", "md", "utf8", "gbk"],
        "md": ["docx", "html"],
        "html": ["txt", "docx"],
        "pptx": ["pdf", "txt"],
        "ppt": ["pdf", "txt"],
    }

    def _mk(fn, label):
        def wrapper(src, tgt, out_dir, progress, opts):
            out = unique_path(out_dir, stem(src), tgt)
            progress.report(0, 100, f"开始：{label}")
            result = fn(src, out, progress)
            return [result] if isinstance(result, str) else list(result)

        return wrapper

    def _pdf_images_wrapper(image_format, label):
        """PDF → 多张图片：直接输出到目录，每页一张"""
        def wrapper(src, tgt, out_dir, progress, opts):
            progress.report(0, 100, f"开始：{label}")
            results = pdf_to_images(
                src, out_dir, dpi=150, image_format=image_format, progress=progress)
            return list(results)

        return wrapper

    registry.register("文档", "pdf", "docx", _mk(pdf_to_docx, "PDF → Word"))
    registry.register("文档", "pdf", "txt", _mk(pdf_to_text, "PDF → 文本"))
    registry.register("文档", "pdf", "png",
                      _pdf_images_wrapper("png", "PDF → 图片"))
    registry.register("文档", "pdf", "jpg",
                      _pdf_images_wrapper("jpg", "PDF → 图片"))
    registry.register("文档", "pdf", "webp",
                      _pdf_images_wrapper("webp", "PDF → 图片"))

    registry.register("文档", "docx", "pdf", _mk(docx_to_pdf, "Word → PDF"))
    registry.register("文档", "docx", "txt", _mk(docx_to_text, "Word → 文本"))
    registry.register("文档", "docx", "md", _mk(docx_to_md, "Word → Markdown"))
    registry.register("文档", "doc", "pdf", _mk(docx_to_pdf, "Word → PDF"))
    registry.register("文档", "doc", "txt", _mk(docx_to_text, "Word → 文本"))
    registry.register("文档", "doc", "docx",
                      _mk(lambda s, o, p: _office_com_convert(s, o, ("Word.Application", "KWps.Application"), 16), "DOC → DOCX"))

    registry.register("文档", "txt", "docx", _mk(text_to_docx, "文本 → Word"))
    registry.register("文档", "txt", "md", _mk(txt_to_md, "文本 → Markdown"))
    registry.register("文档", "md", "docx", _mk(markdown_to_docx, "Markdown → Word"))
    registry.register("文档", "md", "html", _mk(md_to_html, "Markdown → HTML"))
    registry.register("文档", "html", "txt", _mk(html_to_text, "HTML → 文本"))
    registry.register("文档", "html", "docx", _mk(html_to_docx, "HTML → Word"))

    registry.register("文档", "xlsx", "csv", _mk(xlsx_to_csv, "Excel → CSV"))
    registry.register("文档", "xlsx", "txt", _mk(xlsx_to_txt, "Excel → 文本"))
    registry.register("文档", "xls", "xlsx", _mk(xls_to_xlsx, "XLS → XLSX"))
    registry.register("文档", "csv", "xlsx", _mk(csv_to_xlsx, "CSV → Excel"))

    # 文本编码转换（修复乱码）
    for _src in ("txt", "csv"):
        registry.register("文档", _src, "utf8",
                          _encoding_wrapper("utf-8", "转为 UTF-8"))
        registry.register("文档", _src, "gbk",
                          _encoding_wrapper("gbk", "转为 GBK"))

    registry.register("文档", "pptx", "pdf", _mk(pptx_to_pdf, "PPT → PDF"))
    registry.register("文档", "pptx", "txt", _mk(pptx_to_text, "PPT → 文本"))
    registry.register("文档", "ppt", "pdf", _mk(pptx_to_pdf, "PPT → PDF"))
    registry.register("文档", "ppt", "txt", _mk(pptx_to_text, "PPT → 文本"))

    return targets
